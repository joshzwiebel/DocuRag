import os
import io
import logging
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from googleapiclient.errors import HttpError
from dotenv import load_dotenv
import PyPDF2
from io import BytesIO
import docx
import pandas as pd
from bs4 import BeautifulSoup
import pypdf
import xml.etree.ElementTree as ET
import json

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class GoogleDriveClient:
    def __init__(self, service_account_file, scopes):
        self.service_account_file = service_account_file
        self.scopes = scopes
        self.service = self.authenticate()

    def authenticate(self):
        logger.info("Authenticating with service account file: %s", self.service_account_file)
        credentials = service_account.Credentials.from_service_account_file(
            self.service_account_file, scopes=self.scopes)
        return build('drive', 'v3', credentials=credentials)

    def list_files(self, page_size=100):
        logger.info("Listing files with page size: %d", page_size)
        results = self.service.files().list(pageSize=page_size, fields="nextPageToken, files(id, name, mimeType)").execute()
        items = results.get('files', [])

        if not items:
            logger.info('No files found.')
        else:
            logger.info('Files:')
            for item in items:
                logger.info("%s (%s)", item['name'], item['id'])
        return items

    def download_file(self, file_id, file_name):
        logger.info("Downloading file with ID: %s to %s", file_id, file_name)
        request = self.service.files().get_media(fileId=file_id)
        fh = io.FileIO(file_name, 'wb')
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
            logger.info("Download %d%%.", int(status.progress() * 100))

    def get_file_content(self, file_id):
        try:
            # Attempt to download the file content
            request = self.service.files().get_media(fileId=file_id)
            fh = io.BytesIO()
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done:
                status, done = downloader.next_chunk()
                logger.info("Download %d%%.", int(status.progress() * 100))
            fh.seek(0)
            return fh.read().decode('utf-8')
        except HttpError as error:
            if error.resp.status == 403 and 'fileNotDownloadable' in str(error):
                # Handle Google Docs Editors files by exporting them
                logger.info("File is not directly downloadable. Attempting to export.")
                request = self.service.files().export(fileId=file_id, mimeType='text/plain')
                response = request.execute()
                return response.decode('utf-8')
            else:
                logger.error("An error occurred: %s", error)
                raise ValueError(f"An error occurred: {error}")

    def get_pdf_text_content(self, file_id):
        try:
            # Download the PDF file content
            request = self.service.files().get_media(fileId=file_id)
            fh = io.BytesIO()
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done:
                status, done = downloader.next_chunk()
                logger.info("Download %d%%.", int(status.progress() * 100))
            fh.seek(0)
            
            # Extract text from the PDF file
            reader = PyPDF2.PdfFileReader(fh)
            text_content = ""
            for page_num in range(reader.numPages):
                page = reader.getPage(page_num)
                text_content += page.extract_text()
            return text_content
        except HttpError as error:
            logger.error("An error occurred while downloading or processing the PDF file %s: %s", file_id, error)
            return None

    def export_file(self, file_id, mime_type):
        try:
            if mime_type == 'application/vnd.google-apps.document':
                export_mime_type = 'text/plain'
            elif mime_type == 'application/vnd.google-apps.spreadsheet':
                export_mime_type = 'text/csv'
            elif mime_type == 'application/vnd.google-apps.presentation':
                export_mime_type = 'application/pdf'
            else:
                export_mime_type = 'application/pdf'  # Default export type

            request = self.service.files().export(fileId=file_id, mimeType=export_mime_type)
            response = request.execute()
            return response.decode('utf-8')
        except HttpError as error:
            logger.error("An error occurred while exporting file %s: %s", file_id, error)
            return None

    def list_files_in_folder(self, folder_id):
        logger.info("Listing files in folder with ID: %s", folder_id)
        results = self.service.files().list(q=f"'{folder_id}' in parents", fields="files(id, name)").execute()
        items = results.get('files', [])
        if not items:
            logger.info('No files found.')
        else:
            logger.info('Files:')
            for item in items:
                logger.info("%s (%s)", item['name'], item['id'])

    def extract_text(self, file_id, mime_type):
        """Extract text from various file formats."""
        try:
            # Download the file content
            request = self.service.files().get_media(fileId=file_id)
            file_content = io.BytesIO()
            downloader = MediaIoBaseDownload(file_content, request)
            done = False
            while not done:
                status, done = downloader.next_chunk()
            file_content.seek(0)
            
            # Handle different file types
            if mime_type == 'application/pdf':
                return self._extract_pdf(file_content)
            elif mime_type in ['application/msword', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document']:
                return self._extract_word(file_content)
            elif mime_type in ['application/vnd.ms-excel', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet']:
                return self._extract_excel(file_content)
            elif mime_type in ['application/vnd.ms-powerpoint', 'application/vnd.openxmlformats-officedocument.presentationml.presentation']:
                return self._extract_powerpoint(file_content)
            elif mime_type in ['text/plain', 'text/csv', 'text/markdown']:
                return file_content.getvalue().decode('utf-8', errors='ignore')
            elif mime_type == 'application/json':
                return json.loads(file_content.getvalue())
            elif mime_type == 'application/xml':
                return self._extract_xml(file_content)
            elif mime_type == 'text/html':
                return self._extract_html(file_content)

            
            logger.warning(f"Unsupported mime type: {mime_type}")
            return None
            
        except Exception as error:
            logger.error(f"Error extracting text from file {file_id}: {error}")
            return None

    def _extract_pdf(self, file_stream):
        """Extract text from PDF."""
        reader = pypdf.PdfReader(file_stream)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text

    def _extract_word(self, file_stream):
        """Extract text from Word documents."""
        doc = docx.Document(file_stream)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])

    def _extract_excel(self, file_stream):
        """Extract text from Excel files."""
        df = pd.read_excel(file_stream)
        return df.to_string()

    def _extract_powerpoint(self, file_stream):
        """Extract text from PowerPoint files."""
        from pptx import Presentation
        prs = Presentation(file_stream)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def _extract_xml(self, file_stream):
        """Extract text from XML files."""
        tree = ET.parse(file_stream)
        return ET.tostring(tree.getroot(), encoding='unicode', method='text')

    def _extract_html(self, file_stream):
        """Extract text from HTML files."""
        soup = BeautifulSoup(file_stream, 'html.parser')
        return soup.get_text()


if __name__ == '__main__':
    load_dotenv()
    SERVICE_ACCOUNT_FILE = os.getenv("SERVICE_ACCOUNT_FILE")
    SCOPES = os.getenv("SCOPES")
    gdrive_client = GoogleDriveClient(SERVICE_ACCOUNT_FILE, [SCOPES])
    gdrive_client.list_files()